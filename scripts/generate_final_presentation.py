from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets" / "presentation"
REQ = ASSETS / "req"
SCREENS = ASSETS / "screens"
OUT = ROOT / "output" / "ppt"
PPTX = OUT / "junkio-final-2026-20slides.pptx"

IMG = {
    "dashboard": SCREENS / "dashboard.png",
    "transactions": SCREENS / "transactions.png",
    "goals": SCREENS / "goals.png",
    "family": SCREENS / "family.png",
    "admin": SCREENS / "admin.png",
    "swagger": SCREENS / "swagger.png",
    "frontend_req": REQ / "yc2.png",
    "api_req": REQ / "yc3.jpg",
    "devops_req": REQ / "yc4.jpg",
}

BG = RGBColor(246, 248, 252)
CARD = RGBColor(255, 255, 255)
NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(71, 85, 105)
MUTED = RGBColor(100, 116, 139)
BLUE = RGBColor(14, 165, 233)
TEAL = RGBColor(20, 184, 166)
GOLD = RGBColor(245, 158, 11)
GREEN = RGBColor(34, 197, 94)
RED = RGBColor(239, 68, 68)
BORDER = RGBColor(226, 232, 240)
LBLUE = RGBColor(224, 242, 254)
LTEAL = RGBColor(204, 251, 241)
LGOLD = RGBColor(254, 243, 199)
LRED = RGBColor(254, 226, 226)


def I(value):
    return Inches(value)


def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill, line=None, radius=True, alpha=0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = alpha
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def text(
    slide,
    x,
    y,
    w,
    h,
    value,
    size=20,
    color=NAVY,
    bold=False,
    align=PP_ALIGN.LEFT,
    italic=False,
    font="Segoe UI",
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    paragraph.text = value
    run = paragraph.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tb


def bullets(slide, x, y, w, h, items, size=17, bullet_color=BLUE, color=SLATE, gap=5):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r1 = p.add_run()
        r1.text = "• "
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = "Segoe UI"
    return tb


def chip(slide, x, y, w, label, fill, color=NAVY):
    box(slide, x, y, w, I(0.34), fill, radius=True)
    text(slide, x, y + I(0.05), w, I(0.16), label, 10.5, color, True, PP_ALIGN.CENTER)


def header(slide, title_value, sub_value, number):
    box(slide, I(0.65), I(0.42), I(1.05), I(0.12), BLUE, radius=True)
    text(slide, I(0.8), I(0.2), I(4.5), I(0.2), "JUNKIO EXPENSE TRACKER", 11, BLUE, True)
    text(slide, I(0.8), I(0.56), I(10.7), I(0.42), title_value, 26, NAVY, True)
    text(slide, I(0.8), I(1.02), I(11.0), I(0.24), sub_value, 12.5, MUTED)
    text(slide, I(12.15), I(0.3), I(0.45), I(0.18), str(number), 12, MUTED, True, PP_ALIGN.RIGHT)


def fit(slide, path, x, y, w, h, cap=None):
    box(slide, x, y, w, h, CARD, BORDER)
    path = Path(path)
    if not path.exists():
        box(slide, x + I(0.12), y + I(0.12), w - I(0.24), h - I(0.24), LBLUE, BLUE)
        text(slide, x + I(0.15), y + h / 2 - I(0.08), w - I(0.3), I(0.18), "Thêm hình tại đây", 13, NAVY, True, PP_ALIGN.CENTER)
        return
    with Image.open(path) as image:
        iw, ih = image.size
    mx = I(0.12)
    mt = I(0.12)
    mb = I(0.35) if cap else I(0.12)
    bw = w - 2 * mx
    bh = h - mt - mb
    ir = iw / ih
    br = bw / bh
    if ir > br:
        pw = bw
        ph = bw / ir
    else:
        ph = bh
        pw = bh * ir
    px = x + mx + (bw - pw) / 2
    py = y + mt + (bh - ph) / 2
    slide.shapes.add_picture(str(path), int(px), int(py), width=int(pw), height=int(ph))
    if cap:
        text(slide, x + I(0.14), y + h - I(0.24), w - I(0.28), I(0.14), cap, 10.5, MUTED, True, PP_ALIGN.CENTER)


def cover_image(slide, path, sw, sh):
    path = Path(path)
    if not path.exists():
        bg(slide, NAVY)
        return
    with Image.open(path) as image:
        iw, ih = image.size
    ir = iw / ih
    sr = sw / sh
    if ir > sr:
        nh = sh
        nw = sh * ir
        px = -(nw - sw) / 2
        py = 0
    else:
        nw = sw
        nh = sw / ir
        px = 0
        py = -(nh - sh) / 2
    slide.shapes.add_picture(str(path), int(px), int(py), width=int(nw), height=int(nh))


def conn(slide, x1, y1, x2, y2, color=BORDER, width=1.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def stat_card(slide, x, y, w, title_value, value, fill=LBLUE):
    box(slide, x, y, w, I(1.0), CARD, BORDER)
    box(slide, x + I(0.14), y + I(0.16), I(0.1), I(0.68), fill, radius=True)
    text(slide, x + I(0.34), y + I(0.18), w - I(0.5), I(0.16), title_value, 11.5, MUTED, True)
    text(slide, x + I(0.34), y + I(0.44), w - I(0.45), I(0.2), value, 18, NAVY, True)


def section_card(slide, x, y, w, h, title_value, items, tint=LBLUE, bullet_color=BLUE):
    box(slide, x, y, w, h, CARD, BORDER)
    box(slide, x, y, w, I(0.5), tint, radius=True)
    text(slide, x + I(0.2), y + I(0.18), w - I(0.4), I(0.18), title_value, 17, NAVY, True)
    bullets(slide, x + I(0.2), y + I(0.68), w - I(0.4), h - I(0.9), items, 14.2, bullet_color)


def label_row(slide, x, y, labels, fills):
    cursor = x
    for idx, label in enumerate(labels):
        width = 0.55 + len(label) * 0.07
        chip(slide, cursor, y, I(width), label, fills[idx % len(fills)])
        cursor += I(width + 0.12)


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cover_image(slide, IMG["dashboard"], prs.slide_width, prs.slide_height)
    box(slide, 0, 0, prs.slide_width, prs.slide_height, NAVY, radius=False, alpha=0.3)
    box(slide, I(0.7), I(0.72), I(6.4), I(5.95), NAVY, radius=True, alpha=0.1)
    text(slide, I(0.96), I(0.95), I(3.3), I(0.2), "BÁO CÁO THUYẾT TRÌNH CUỐI KỲ", 12, LGOLD, True)
    text(slide, I(0.96), I(1.48), I(6.0), I(0.78), "Junkio Expense Tracker", 28, CARD, True)
    text(slide, I(0.96), I(2.35), I(5.7), I(0.36), "Ứng dụng web quản lý chi tiêu cá nhân và gia đình", 18, RGBColor(226, 232, 240))
    text(slide, I(0.96), I(3.18), I(4.9), I(0.2), "Sinh viên: [Điền họ tên]", 15, CARD)
    text(slide, I(0.96), I(3.56), I(4.9), I(0.2), "MSSV: [Điền MSSV]", 15, CARD)
    text(slide, I(0.96), I(3.94), I(4.9), I(0.2), "Lớp: [Điền lớp]", 15, CARD)
    text(slide, I(0.96), I(4.32), I(4.9), I(0.2), "GVHD: [Điền tên giảng viên]", 15, CARD)
    label_row(slide, I(0.96), I(5.15), ["React + Vite", "Node.js + Express", "Docker Compose"], [LBLUE, LTEAL, LGOLD])
    text(slide, I(0.96), I(5.82), I(5.95), I(0.3), "Mục tiêu của phần trình bày là chứng minh hệ thống đã được hoàn thiện ở cả kiến trúc, chức năng, bảo mật, kiểm thử và đóng gói triển khai.", 12.2, RGBColor(226, 232, 240))


def add_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Nội dung thuyết trình", "Bản mở rộng 20 slide đi theo thứ tự từ bài toán, kiến trúc, demo chức năng đến đánh giá hoàn thiện.", 2)
    topics = [
        ("01", "Bài toán và mục tiêu", "Nêu nhu cầu thực tế, mục tiêu và phạm vi sản phẩm."),
        ("02", "Kiến trúc và công nghệ", "Giải thích FE/BE, dữ liệu, bảo mật, Docker."),
        ("03", "Chức năng và luồng demo", "Đi qua các nghiệp vụ chính, API, dashboard, admin."),
        ("04", "Đánh giá và kết luận", "Kiểm thử, mức hoàn thiện, hạn chế và hướng phát triển."),
    ]
    y = 1.8
    fills = [LBLUE, LTEAL, LGOLD, LRED]
    for idx, (num, title_value, desc) in enumerate(topics):
        box(slide, I(1.0), I(y), I(11.25), I(0.95), CARD, BORDER)
        box(slide, I(1.18), I(y + 0.17), I(0.68), I(0.56), fills[idx], radius=True)
        text(slide, I(1.18), I(y + 0.31), I(0.68), I(0.16), num, 14, NAVY, True, PP_ALIGN.CENTER)
        text(slide, I(2.12), I(y + 0.18), I(3.6), I(0.18), title_value, 18, NAVY, True)
        text(slide, I(2.12), I(y + 0.47), I(8.8), I(0.16), desc, 12.5, SLATE)
        y += 1.12
    box(slide, I(1.0), I(6.4), I(11.25), I(0.42), NAVY)
    text(slide, I(1.25), I(6.51), I(10.8), I(0.16), "Kịch bản này phù hợp cho thời lượng báo cáo 12-15 phút nhưng vẫn đủ chiều sâu khi giảng viên hỏi kỹ về kiến trúc hoặc DevOps.", 11.8, CARD)


def add_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Bài toán và lý do chọn đề tài", "Đề tài xuất phát từ nhu cầu quản lý tài chính có hệ thống và có khả năng mở rộng sang bối cảnh gia đình.", 3)
    section_card(slide, I(0.9), I(1.72), I(5.65), I(4.9), "Vấn đề thực tế", [
        "Chi tiêu cá nhân thường được lưu rời rạc ở sổ tay, ghi chú hoặc nhiều ứng dụng nhỏ lẻ.",
        "Khi phát sinh quỹ chung gia đình, việc theo dõi thành viên đóng góp và khoản chi dùng chung trở nên khó minh bạch.",
        "Người dùng cần không chỉ CRUD giao dịch mà còn cần dashboard, mục tiêu, ngân sách và cảnh báo sớm.",
        "Đây là bài toán phù hợp để thể hiện năng lực full-stack và tư duy thiết kế hệ thống web thực tế.",
    ], LBLUE, BLUE)
    box(slide, I(6.85), I(1.72), I(5.4), I(4.9), NAVY)
    text(slide, I(7.18), I(2.02), I(4.7), I(0.2), "Thông điệp chính", 16, LGOLD, True)
    text(slide, I(7.18), I(2.48), I(4.45), I(1.56), "Junkio không chỉ là ứng dụng ghi chép thu chi, mà là một nền tảng quản lý tài chính với kiến trúc SPA, API chuẩn hóa, phân quyền, báo cáo và khả năng demo triển khai thực tế.", 21, CARD, True)
    label_row(slide, I(7.18), I(5.4), ["SPA", "REST API", "Realtime", "Docker"], [LBLUE, LTEAL, LGOLD, LRED])


def add_objectives(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Mục tiêu và phạm vi hệ thống", "Slide này giúp hội đồng thấy rõ đề tài đang làm gì, làm đến đâu và phạm vi kỹ thuật được bao phủ.", 4)
    section_card(slide, I(0.9), I(1.75), I(4.0), I(4.75), "Mục tiêu chính", [
        "Tách biệt frontend và backend theo mô hình client-server.",
        "Quản lý giao dịch, ví, ngân sách, mục tiêu, gia đình và thông báo.",
        "Đảm bảo xác thực, phân quyền và validation dữ liệu.",
        "Có dashboard để trực quan hóa số liệu tài chính.",
    ], LBLUE, BLUE)
    section_card(slide, I(4.98), I(1.75), I(3.55), I(4.75), "Phạm vi kỹ thuật", [
        "Frontend SPA với routing phía client.",
        "Backend RESTful API trả JSON.",
        "PostgreSQL làm hệ quản trị dữ liệu chính.",
        "Redis và Socket.IO phục vụ cache/realtime.",
    ], LTEAL, TEAL)
    section_card(slide, I(8.62), I(1.75), I(3.55), I(4.75), "Tiêu chí hoàn thiện", [
        "Chạy ổn định khi demo.",
        "Có Docker Compose để dựng toàn hệ thống.",
        "Có test và tài liệu API để minh chứng.",
        "Có thể trình diễn được luồng nghiệp vụ thật.",
    ], LGOLD, GOLD)


def add_requirements(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Đối chiếu với yêu cầu môn học", "Ba mốc báo cáo đều đã được gom lại thành các nhóm năng lực chính: frontend, API, integration và DevOps.", 5)
    box(slide, I(0.86), I(1.7), I(3.45), I(4.95), CARD, BORDER)
    text(slide, I(1.1), I(1.96), I(2.5), I(0.2), "Nhóm yêu cầu chính", 18, NAVY, True)
    bullets(slide, I(1.1), I(2.34), I(2.8), I(4.0), [
        "Frontend theo hướng SPA, responsive và chia component tái sử dụng.",
        "Backend có RESTful API, JWT, validation và RBAC.",
        "Có tích hợp từ giao diện web đến database thực tế.",
        "Có Docker, Swagger/Postman và báo cáo dashboard.",
        "Có thể trình diễn 'one click' với docker compose.",
    ], 15.2, GREEN)
    fit(slide, IMG["frontend_req"], I(4.56), I(1.7), I(3.85), I(2.35), "Mốc frontend và component")
    fit(slide, IMG["api_req"], I(8.62), I(1.7), I(3.55), I(2.35), "Mốc API, Auth, Integration")
    fit(slide, IMG["devops_req"], I(4.56), I(4.28), I(7.61), I(2.37), "Mốc Dockerization và báo cáo tổng kết")


def add_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Kiến trúc tổng thể hệ thống", "Kiến trúc được triển khai theo mô hình nhiều service, mỗi thành phần có trách nhiệm rõ ràng.", 6)
    nodes = [
        ("Browser", 0.95, 2.35, 1.55, 0.72, LBLUE),
        ("Frontend\nReact + Vite", 2.9, 1.9, 2.2, 0.95, LTEAL),
        ("Nginx\nStatic + Proxy", 2.9, 3.3, 2.2, 0.95, LGOLD),
        ("Backend API\nNode.js + Express", 5.95, 2.55, 2.4, 1.05, CARD),
        ("PostgreSQL", 9.3, 1.86, 2.15, 0.85, LBLUE),
        ("Redis", 9.3, 3.22, 2.15, 0.85, LTEAL),
        ("Swagger / Postman", 9.08, 4.62, 2.58, 0.85, LGOLD),
    ]
    for label, x, y, w, h, fill in nodes:
        box(slide, I(x), I(y), I(w), I(h), fill, BORDER)
        text(slide, I(x + 0.06), I(y + 0.2), I(w - 0.12), I(h - 0.24), label, 15.5, NAVY, True, PP_ALIGN.CENTER)
    conn(slide, I(2.5), I(2.72), I(2.9), I(2.38), BLUE)
    conn(slide, I(2.5), I(2.72), I(2.9), I(3.78), BLUE)
    conn(slide, I(5.1), I(2.38), I(5.95), I(2.98), TEAL)
    conn(slide, I(5.1), I(3.78), I(5.95), I(3.18), TEAL)
    conn(slide, I(8.35), I(2.98), I(9.3), I(2.28), BLUE)
    conn(slide, I(8.35), I(3.2), I(9.3), I(3.64), TEAL)
    conn(slide, I(8.35), I(3.2), I(9.08), I(5.04), GOLD)
    box(slide, I(0.95), I(5.78), I(10.75), I(0.58), CARD, BORDER)
    text(slide, I(1.18), I(5.94), I(10.25), I(0.18), "Luồng chính: Browser → Frontend/Nginx → Backend API → PostgreSQL và Redis. Mô hình này bám sát cách triển khai thật khi dựng bằng Docker Compose.", 12.6, SLATE)
    chip(slide, I(11.9), I(5.9), I(0.85), "Docker", LGOLD)


def add_layers(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Kiến trúc tầng và trách nhiệm từng khối", "Không chỉ tách service, hệ thống còn tách rõ tầng giao diện, tầng nghiệp vụ, tầng dữ liệu và tầng hạ tầng.", 7)
    cards = [
        ("Tầng giao diện", ["Trang dashboard, transactions, goals, family, admin.", "Nhận input người dùng, validate sơ bộ và gọi API."], LBLUE, BLUE, 0.9),
        ("Tầng nghiệp vụ", ["Controller, service, middleware và auth flow.", "Xử lý rule, phân quyền, validation và format response."], LTEAL, TEAL, 4.35),
        ("Tầng dữ liệu", ["Sequelize model, migration, seeder, PostgreSQL.", "Lưu trữ giao dịch, ví, người dùng, gia đình, mục tiêu."], LGOLD, GOLD, 7.8),
    ]
    for title_value, items, fill, accent, xpos in cards:
        section_card(slide, I(xpos), I(1.85), I(3.05), I(4.75), title_value, items, fill, accent)
    box(slide, I(0.9), I(6.0), I(10.8), I(0.42), NAVY)
    text(slide, I(1.12), I(6.11), I(10.3), I(0.16), "Lợi ích: dễ mở rộng, dễ kiểm thử, giảm coupling giữa UI, business logic và data layer.", 11.8, CARD)
    fit(slide, IMG["admin"], I(11.0), I(1.85), I(1.95), I(2.7), "Admin")


def add_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Công nghệ sử dụng", "Stack được chọn để cân bằng giữa tốc độ phát triển, khả năng mở rộng và tính phù hợp với đề tài môn học.", 8)
    columns = [
        ("Frontend", ["React 19", "Vite", "Tailwind CSS", "Redux Toolkit", "React Router", "Formik + Yup", "Recharts"], LBLUE, BLUE, 0.9),
        ("Backend", ["Node.js", "Express", "Sequelize ORM", "PostgreSQL", "JWT", "Socket.IO", "Redis", "Swagger"], LTEAL, TEAL, 4.55),
        ("QA / DevOps", ["Docker", "docker-compose", "Jest", "Vitest", "Playwright", "Newman", "Git / GitHub"], LGOLD, GOLD, 8.2),
    ]
    for title_value, items, fill, accent, xpos in columns:
        box(slide, I(xpos), I(1.75), I(3.1), I(4.85), CARD, BORDER)
        box(slide, I(xpos), I(1.75), I(3.1), I(0.5), fill, radius=True)
        text(slide, I(xpos + 0.2), I(1.95), I(2.3), I(0.16), title_value, 18, NAVY, True)
        bullets(slide, I(xpos + 0.22), I(2.48), I(2.65), I(3.8), items, 14.4, accent)
    text(slide, I(1.0), I(6.78), I(11.4), I(0.18), "Thông điệp chính: công nghệ không chỉ để 'làm được' mà còn để thể hiện rõ tính hiện đại, bảo mật và khả năng demo ổn định.", 12.5, MUTED, False, PP_ALIGN.CENTER, True)


def add_structure(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Cấu trúc source code và tư duy tổ chức dự án", "Repo được tách thành frontend, backend, docker và tài liệu; cách tổ chức này giúp trình bày và bảo trì thuận lợi.", 9)
    box(slide, I(0.9), I(1.8), I(5.2), I(4.75), CARD, BORDER)
    text(slide, I(1.12), I(2.02), I(3.0), I(0.18), "Phân rã theo thư mục", 18, NAVY, True)
    tree_lines = [
        "frontend/",
        "  src/pages, src/components, src/hooks, src/store",
        "backend/",
        "  controllers, routes, models, middleware, config, tests",
        "docker-compose.yml",
        "scripts/, doc/, assets/presentation/",
    ]
    y = 2.42
    for line in tree_lines:
        text(slide, I(1.18), I(y), I(4.3), I(0.18), line, 14.8, SLATE, False, font="Consolas")
        y += 0.42
    section_card(slide, I(6.35), I(1.8), I(2.7), I(2.2), "Ý nghĩa", [
        "Dễ tách phần FE/BE để demo.",
        "Dễ quy trách nhiệm từng khối.",
        "Thuận tiện cho test và Docker.",
    ], LBLUE, BLUE)
    section_card(slide, I(9.25), I(1.8), I(2.95), I(2.2), "Tư duy component", [
        "Page là composition của nhiều component nhỏ.",
        "State dùng cục bộ hoặc Redux tùy cấp độ dùng lại.",
        "Có thể mở rộng mà không phá vỡ toàn cục.",
    ], LTEAL, TEAL)
    section_card(slide, I(6.35), I(4.25), I(5.85), I(2.3), "Thông điệp khi trình bày", [
        "Không chỉ có giao diện chạy được, dự án còn có cấu trúc code đủ sạch để người khác tiếp nhận và phát triển tiếp.",
        "Việc đưa toàn bộ asset thuyết trình vào repo cũng giúp giảm rủi ro path dài và tăng tính tái lập khi tạo slide.",
    ], LGOLD, GOLD)


def add_data_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Thiết kế dữ liệu và nghiệp vụ lõi", "Dữ liệu xoay quanh transaction, nhưng được mở rộng thêm để đáp ứng ví, mục tiêu, ngân sách và hoạt động gia đình.", 10)
    box(slide, I(5.1), I(2.42), I(2.25), I(0.95), BLUE)
    text(slide, I(5.35), I(2.74), I(1.75), I(0.18), "Transactions", 22, CARD, True, PP_ALIGN.CENTER)
    entities = [
        ("Users", 1.05, 1.42, LBLUE),
        ("Wallets", 1.05, 4.08, LTEAL),
        ("Categories", 4.72, 1.02, LGOLD),
        ("Budgets", 8.72, 1.02, LBLUE),
        ("Goals", 8.98, 4.15, LTEAL),
        ("Families", 11.02, 2.36, LGOLD),
        ("Notifications", 4.68, 5.45, LBLUE),
    ]
    for name, xpos, ypos, fill in entities:
        box(slide, I(xpos), I(ypos), I(1.85), I(0.72), fill, BORDER)
        text(slide, I(xpos), I(ypos + 0.22), I(1.85), I(0.18), name, 14.2, NAVY, True, PP_ALIGN.CENTER)
        conn(slide, I(xpos + 0.92), I(ypos + 0.36), I(6.22), I(2.92), BORDER)
    box(slide, I(0.92), I(5.95), I(12.0), I(0.48), CARD, BORDER)
    text(slide, I(1.15), I(6.08), I(11.45), I(0.16), "Migrations dùng để version hóa schema; seeder/demo data hỗ trợ khởi tạo nhanh dữ liệu cho quá trình test, demo và bảo vệ.", 12.2, SLATE)


def add_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Bảo mật, xác thực và phân quyền", "Đây là phần quan trọng để chứng minh hệ thống không chỉ có giao diện mà còn có năng lực xử lý dữ liệu an toàn.", 11)
    section_card(slide, I(0.9), I(1.82), I(3.4), I(4.65), "Xác thực", [
        "Đăng nhập tạo JWT access token.",
        "Refresh token hỗ trợ duy trì phiên.",
        "Quên mật khẩu và reset password đã được triển khai.",
    ], LBLUE, BLUE)
    section_card(slide, I(4.55), I(1.82), I(3.4), I(4.65), "Bảo vệ dữ liệu", [
        "Password được hash bằng bcrypt.",
        "API dùng middleware auth trước các route protected.",
        "CORS, Helmet, rate-limit giúp giảm rủi ro phổ biến.",
    ], LTEAL, TEAL)
    section_card(slide, I(8.2), I(1.82), I(4.0), I(4.65), "Phân quyền", [
        "Role chính gồm admin, staff, member.",
        "Mỗi role chỉ được truy cập tập API phù hợp.",
        "Admin dashboard và luồng quản trị là minh chứng rõ nhất cho RBAC.",
    ], LGOLD, GOLD)
    label_row(slide, I(1.0), I(6.18), ["Admin", "Staff", "Member", "JWT", "Hash"], [LGOLD, LTEAL, LBLUE, LBLUE, LRED])


def add_frontend_components(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Frontend: UI/UX và tư duy component", "Kiến trúc frontend bám đúng mục tiêu của mốc báo cáo tuần 10: SPA, component hóa, responsive và routing mượt.", 12)
    section_card(slide, I(0.9), I(1.78), I(3.5), I(4.95), "Điểm nhấn frontend", [
        "SPA bằng React, chuyển trang không reload.",
        "Redux Toolkit quản lý state và async flow.",
        "Formik/Yup hỗ trợ form và validate phía client.",
        "Tailwind CSS giúp xây UI nhanh và đồng nhất.",
        "Component tái sử dụng giữa nhiều màn hình.",
    ], LBLUE, BLUE)
    section_card(slide, I(4.65), I(1.78), I(3.65), I(2.3), "Ví dụ component", [
        "MainLayout, Navbar, Sidebar, Modal",
        "PageHeader, Card, Table, Form Input",
        "Chart wrapper, filter bar, responsive grid",
    ], LTEAL, TEAL)
    section_card(slide, I(4.65), I(4.42), I(3.65), I(2.3), "Responsive", [
        "Bố cục thích ứng desktop, tablet, mobile.",
        "Giảm mật độ thông tin ở màn hình nhỏ.",
        "Không phá flow thao tác khi đổi kích thước.",
    ], LGOLD, GOLD)
    fit(slide, IMG["dashboard"], I(8.55), I(1.78), I(3.7), I(4.95), "Dashboard thực tế")


def add_frontend_screens(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Các màn hình chính trên frontend", "Slide này rất phù hợp để nói nhanh về những phần mà hội đồng nhìn thấy trực tiếp trên trình duyệt.", 13)
    fit(slide, IMG["dashboard"], I(0.86), I(1.72), I(3.5), I(2.25), "Dashboard")
    fit(slide, IMG["transactions"], I(4.55), I(1.72), I(3.5), I(2.25), "Transactions")
    fit(slide, IMG["goals"], I(8.24), I(1.72), I(3.95), I(2.25), "Goals")
    fit(slide, IMG["family"], I(0.86), I(4.2), I(3.5), I(2.25), "Family / Shared Expense")
    fit(slide, IMG["admin"], I(4.55), I(4.2), I(3.5), I(2.25), "Admin Dashboard")
    box(slide, I(8.24), I(4.2), I(3.95), I(2.25), CARD, BORDER)
    bullets(slide, I(8.46), I(4.46), I(3.45), I(1.55), [
        "Các màn hình đã đủ để thể hiện luồng người dùng và luồng quản trị.",
        "Có thể lấy ngay các màn hình này làm minh họa chính trong buổi demo.",
        "Ảnh đã được nhúng vào PPT nên không phụ thuộc file ngoài khi mở slide.",
    ], 14.2, GREEN)


def add_backend_api(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Backend: RESTful API và tài liệu kỹ thuật", "Backend trả JSON chuẩn, có Swagger để mở trực tiếp trong lúc bảo vệ.", 14)
    box(slide, I(0.9), I(1.78), I(4.15), I(4.95), CARD, BORDER)
    text(slide, I(1.12), I(2.02), I(2.8), I(0.18), "Các nhóm API chính", 18, NAVY, True)
    bullets(slide, I(1.12), I(2.38), I(3.55), I(3.8), [
        "/api/auth",
        "/api/transactions",
        "/api/wallets",
        "/api/budgets",
        "/api/goals",
        "/api/families",
        "/api/admin/*",
    ], 15.4, TEAL)
    fit(slide, IMG["swagger"], I(5.3), I(1.78), I(6.95), I(4.95), "Swagger UI phục vụ minh chứng input, output và thử API nhanh")


def add_integration_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Luồng tích hợp từ frontend đến database", "Đây là slide giúp nối frontend và backend lại thành một câu chuyện hoàn chỉnh trong buổi báo cáo.", 15)
    steps = [
        ("1", "Người dùng thao tác trên giao diện", "Ví dụ thêm giao dịch, nạp goal hoặc tạo chi tiêu chung."),
        ("2", "Frontend gọi API", "Request được gửi đến backend qua endpoint phù hợp."),
        ("3", "Backend kiểm tra", "Auth, role và validation được chạy trước khi ghi dữ liệu."),
        ("4", "Database cập nhật", "PostgreSQL lưu dữ liệu; dashboard và danh sách được làm mới."),
    ]
    y = 1.86
    fills = [LBLUE, LTEAL, LGOLD, LRED]
    for idx, (num, title_value, desc) in enumerate(steps):
        box(slide, I(0.92), I(y), I(11.32), I(0.94), CARD, BORDER)
        box(slide, I(1.1), I(y + 0.17), I(0.58), I(0.56), fills[idx], radius=True)
        text(slide, I(1.1), I(y + 0.32), I(0.58), I(0.14), num, 13, NAVY, True, PP_ALIGN.CENTER)
        text(slide, I(1.98), I(y + 0.18), I(3.5), I(0.18), title_value, 16.5, NAVY, True)
        text(slide, I(1.98), I(y + 0.48), I(8.7), I(0.16), desc, 12.3, SLATE)
        y += 1.08
    box(slide, I(0.92), I(6.34), I(11.32), I(0.36), NAVY)
    text(slide, I(1.15), I(6.43), I(10.8), I(0.14), "Nếu hội đồng yêu cầu demo tích hợp, nên dùng đúng flow này để trình bày: Web → API → DB → UI phản hồi.", 11.5, CARD)


def add_business_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Chức năng nổi bật và giá trị nghiệp vụ", "Điểm mạnh của đề tài là không dừng ở CRUD mà có thêm nhiều lớp nghiệp vụ phục vụ tình huống sử dụng thật.", 16)
    section_card(slide, I(0.9), I(1.84), I(2.8), I(4.7), "Giao dịch & ví", [
        "CRUD giao dịch thu chi.",
        "Quản lý nhiều ví.",
        "Filter, search, phân trang.",
    ], LBLUE, BLUE)
    section_card(slide, I(3.95), I(1.84), I(2.8), I(4.7), "Ngân sách & goal", [
        "Theo dõi budget theo danh mục.",
        "Tạo mục tiêu tiết kiệm.",
        "Nạp tiền và xem tiến độ.",
    ], LTEAL, TEAL)
    section_card(slide, I(7.0), I(1.84), I(2.8), I(4.7), "Gia đình", [
        "Nhóm thành viên.",
        "Chi tiêu dùng chung.",
        "Công nợ nội bộ và minh bạch quỹ.",
    ], LGOLD, GOLD)
    section_card(slide, I(10.05), I(1.84), I(2.2), I(4.7), "Nâng cao", [
        "Recurring transaction.",
        "Realtime notification.",
        "Debt simplification.",
    ], LRED, RED)


def add_analytics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Dashboard, báo cáo và phân tích", "Đây là phần tạo chiều sâu cho đồ án vì biến dữ liệu giao dịch thành thông tin có ích cho người dùng.", 17)
    fit(slide, IMG["dashboard"], I(0.92), I(1.82), I(5.35), I(4.75), "Dashboard tổng quan")
    section_card(slide, I(6.55), I(1.82), I(2.65), I(2.15), "Chỉ số quan sát", [
        "Tổng thu, tổng chi, số dư.",
        "Theo dõi theo thời gian.",
        "Trực quan bằng chart/card.",
    ], LBLUE, BLUE)
    section_card(slide, I(9.45), I(1.82), I(2.8), I(2.15), "Báo cáo", [
        "Export CSV/Excel/PDF.",
        "Phù hợp cho thao tác đối chiếu.",
        "Thuận tiện cho trình diễn.",
    ], LGOLD, GOLD)
    section_card(slide, I(6.55), I(4.42), I(5.7), I(2.15), "Ý nghĩa khi báo cáo", [
        "Dashboard là bằng chứng rõ nhất cho việc dữ liệu đã được tổng hợp đúng từ backend lên giao diện.",
        "Nếu muốn gây ấn tượng trong phần demo, nên mở dashboard trước vì đây là màn hình nhìn ra được chiều sâu của toàn hệ thống.",
    ], LTEAL, TEAL)


def add_testing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Kiểm thử, Swagger và mức độ tin cậy", "Hệ thống đã có các lớp kiểm chứng để tránh cảm giác đây chỉ là một prototype giao diện.", 18)
    stat_card(slide, I(0.95), I(1.82), I(2.75), "Backend test", "20 suite / 112 test", LBLUE)
    stat_card(slide, I(3.95), I(1.82), I(2.45), "Smoke API", "Auth + Transaction", LTEAL)
    stat_card(slide, I(6.65), I(1.82), I(2.45), "API docs", "Swagger + Postman", LGOLD)
    stat_card(slide, I(9.35), I(1.82), I(2.7), "Frontend QA", "Vitest + Playwright", LRED)
    section_card(slide, I(0.95), I(3.2), I(5.45), I(3.35), "Những gì có thể demo", [
        "Đăng nhập để nhận JWT và gọi API protected.",
        "Nhập sai dữ liệu để chứng minh validation trả lỗi chuẩn.",
        "Thực hiện CRUD rồi đối chiếu dữ liệu trên dashboard hoặc danh sách.",
        "Mở Swagger để cho thấy tài liệu API có thể dùng ngay.",
    ], LBLUE, GREEN)
    box(slide, I(6.72), I(3.2), I(5.3), I(3.35), NAVY)
    text(slide, I(6.98), I(3.46), I(4.7), I(0.18), "Nhận định QA", 18, CARD, True)
    text(slide, I(6.98), I(3.9), I(4.55), I(1.55), "Dự án đã vượt mức 'chạy được' nhờ có test backend thực thi, tài liệu API, kiểm chứng bằng Docker và dữ liệu demo có thể dùng cho nhiều luồng bảo vệ khác nhau.", 19, RGBColor(226, 232, 240), True)


def add_devops(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Docker hóa và triển khai one-click", "Đây là phần rất quan trọng để chứng minh đề tài không phụ thuộc vào một môi trường cài đặt thủ công.", 19)
    box(slide, I(0.95), I(1.82), I(5.75), I(2.48), NAVY)
    text(slide, I(1.2), I(2.08), I(3.3), I(0.18), "Lệnh demo chính", 16, LGOLD, True)
    text(slide, I(1.2), I(2.45), I(5.0), I(0.22), "docker compose up -d --build", 22, CARD, True, font="Consolas")
    text(slide, I(1.2), I(2.95), I(5.05), I(0.72), "Sau khi dựng xong có thể truy cập:\n• Web: http://localhost\n• API: http://localhost:5000\n• Swagger: http://localhost:5000/api-docs", 13, RGBColor(226, 232, 240), font="Consolas")
    services = [
        ("web", "Nginx + React build", "Port 80", LBLUE, 7.0, 1.82),
        ("api", "Node.js + Express", "Port 5000", LTEAL, 9.78, 1.82),
        ("db", "PostgreSQL 15", "Port 5432", LGOLD, 7.0, 3.12),
        ("redis", "Redis 7", "Port 6379", LRED, 9.78, 3.12),
    ]
    for name, desc, port, fill, xpos, ypos in services:
        box(slide, I(xpos), I(ypos), I(2.35), I(1.02), fill, BORDER)
        text(slide, I(xpos + 0.18), I(ypos + 0.16), I(1.0), I(0.16), name, 16, NAVY, True)
        text(slide, I(xpos + 0.18), I(ypos + 0.43), I(1.85), I(0.14), desc, 11.2, SLATE)
        text(slide, I(xpos + 0.18), I(ypos + 0.68), I(1.45), I(0.14), port, 11.2, MUTED, True)
    section_card(slide, I(0.95), I(4.72), I(11.18), I(1.82), "Điểm mạnh khi bảo vệ", [
        "Container API có thể tự migrate và seed dữ liệu demo khi khởi động.",
        "Các service có healthcheck và phụ thuộc được điều phối trong docker-compose.",
        "Việc dựng lại hệ thống bằng một lệnh là bằng chứng mạnh nhất cho tính tái lập của dự án.",
    ], LGOLD, TEAL)


def add_readiness(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Mức độ hoàn thiện và các điểm đã sẵn sàng", "Slide này nên dùng để trả lời câu hỏi: dự án đã hoàn thiện đến đâu tính đến thời điểm bảo vệ.", 20)
    section_card(slide, I(0.92), I(1.85), I(3.62), I(4.7), "Đã hoàn thiện khá chắc", [
        "FE/BE tách biệt rõ, có SPA và REST API.",
        "JWT, RBAC, dashboard, Swagger, Docker đều đã có.",
        "Nhiều luồng chức năng có thể demo trực tiếp.",
        "Đã có test backend thực thi và môi trường demo thật.",
    ], LBLUE, BLUE)
    section_card(slide, I(4.72), I(1.85), I(3.62), I(4.7), "Các điểm cần nói trung thực", [
        "Seeder dữ liệu lớn 500-1000 bản ghi vẫn là hạng mục nên hoàn thiện thêm.",
        "Secret nên chuyển sạch hơn sang file .env khi bàn giao chính thức.",
        "Một số chi tiết UI/test frontend còn có thể tiếp tục nâng chất lượng.",
    ], LGOLD, GOLD)
    box(slide, I(8.52), I(1.85), I(3.6), I(4.7), NAVY)
    text(slide, I(8.8), I(2.12), I(3.0), I(0.18), "Kết luận đánh giá", 18, CARD, True)
    text(slide, I(8.8), I(2.58), I(2.8), I(1.7), "Nếu loại trừ phần chuẩn bị tài liệu demo, hệ thống đã đủ mạnh để bảo vệ kỹ thuật và thể hiện rõ năng lực xây dựng một ứng dụng web full-stack có thể triển khai.", 20, RGBColor(226, 232, 240), True)
    chip(slide, I(8.8), I(5.4), I(1.05), "Đủ demo", LGOLD)
    chip(slide, I(9.98), I(5.4), I(1.0), "Có chiều sâu", LTEAL)


def add_limitations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "Hạn chế hiện tại và hướng phát triển", "Một slide cân bằng giúp phần kết luận thuyết phục hơn vì thừa nhận đúng các điểm còn lại và đưa ra lộ trình tiếp theo.", 21)
    section_card(slide, I(0.95), I(1.88), I(3.65), I(4.55), "Hạn chế", [
        "Chưa đẩy mạnh dữ liệu seed lên mức 500-1000 bản ghi như yêu cầu mạnh nhất của đề.",
        "Chưa tối ưu hết edge case giao diện ở mọi độ phân giải và trạng thái dữ liệu.",
        "DevOps mới ở mức Docker local, chưa trình bày pipeline CI/CD hoàn chỉnh.",
    ], LRED, RED)
    section_card(slide, I(4.85), I(1.88), I(3.65), I(4.55), "Hướng nâng cấp ngắn hạn", [
        "Bổ sung generator dữ liệu lớn và dashboard stress test.",
        "Đưa toàn bộ secret sang .env và chuẩn hóa cấu hình môi trường.",
        "Tăng coverage test frontend, đặc biệt cho flow form và table.",
    ], LBLUE, BLUE)
    section_card(slide, I(8.75), I(1.88), I(3.4), I(4.55), "Hướng phát triển dài hạn", [
        "Mobile app hoặc PWA offline.",
        "Phân tích chi tiêu thông minh bằng AI/ML.",
        "Kết nối e-wallet hoặc open banking trong tương lai.",
    ], LTEAL, TEAL)


def add_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box(slide, 0, 0, prs.slide_width, prs.slide_height, NAVY, radius=False)
    box(slide, I(0.78), I(0.78), I(11.82), I(5.95), NAVY, BLUE)
    text(slide, I(1.18), I(1.18), I(2.8), I(0.18), "KẾT LUẬN", 12, LGOLD, True)
    text(slide, I(1.18), I(1.98), I(9.7), I(0.72), "Junkio Expense Tracker cho thấy khả năng xây dựng một hệ thống web full-stack có tính thực tiễn, có tổ chức, có bảo mật và có thể triển khai lại bằng Docker.", 27, CARD, True)
    text(slide, I(1.18), I(3.42), I(8.75), I(0.46), "Em xin cảm ơn thầy/cô và hội đồng đã lắng nghe.\nRất mong nhận được câu hỏi và góp ý để hoàn thiện thêm sản phẩm.", 18, RGBColor(226, 232, 240))
    label_row(slide, I(1.18), I(5.22), ["Q&A", "Cảm ơn", "Junkio"], [LGOLD, LTEAL, LBLUE])
    text(slide, I(12.12), I(0.3), I(0.45), I(0.18), "22", 12, RGBColor(203, 213, 225), True, PP_ALIGN.RIGHT)


def add_deck():
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = I(13.333)
    prs.slide_height = I(7.5)
    prs.core_properties.title = "Junkio Expense Tracker - Bao cao cuoi ky"
    prs.core_properties.subject = "Mau PowerPoint 20 slide"
    prs.core_properties.author = "OpenAI Codex"

    add_cover(prs)
    add_agenda(prs)
    add_problem(prs)
    add_objectives(prs)
    add_requirements(prs)
    add_architecture(prs)
    add_layers(prs)
    add_tech_stack(prs)
    add_structure(prs)
    add_data_model(prs)
    add_security(prs)
    add_frontend_components(prs)
    add_frontend_screens(prs)
    add_backend_api(prs)
    add_integration_flow(prs)
    add_business_features(prs)
    add_analytics(prs)
    add_testing(prs)
    add_devops(prs)
    add_readiness(prs)
    add_limitations(prs)
    add_closing(prs)

    prs.save(str(PPTX))
    print(PPTX)


if __name__ == "__main__":
    add_deck()
